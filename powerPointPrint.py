from pptx import Presentation
from pptx.util import Inches
import json
import os
import sys
import codecs
import json

argv = sys.argv
try:
    index = argv.index("--") + 1
except:
    index = len(argv)

argv = argv[index:]
if(len(argv) > 0 and argv[0]):
	output_file = (argv[0])


prs = Presentation()
for i in range(1, len(argv)):
    img_path = argv[i]

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    pic = slide.shapes.add_picture(img_path, 0, 0)

prs.save(output_file)